from pptx import Presentation

def parse_pptx(file_path: str) -> str:
       """Parse text from a PPTX file."""
       try:
           prs = Presentation(file_path)
           text = []
           for slide in prs.slides:
               for shape in slide.shapes:
                   if hasattr(shape, 'text') and shape.text.strip():
                       text.append(shape.text)
           return '\n'.join(text) if text else "No text extracted from PPTX."
       except Exception as e:
           return f"Error parsing PPTX: {str(e)}"