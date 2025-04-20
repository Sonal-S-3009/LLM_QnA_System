import os
import pdfplumber
import docx
import pptx
import pandas as pd
import pytesseract
from PIL import Image
import requests
from bs4 import BeautifulSoup
import re
import json

def extract_text_from_image(file_path):
    """Extract text from images using Tesseract OCR."""
    try:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        return text if text.strip() else "No text found in image."
    except Exception as e:
        return f"OCR Error: {str(e)}"

def extract_text_from_pdf(file_path):
    """Extract text from PDFs, including image-based ones."""
    try:
        with pdfplumber.open(file_path) as pdf:
            text = ""
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
                else:
                    img = page.to_image().original
                    ocr_text = pytesseract.image_to_string(img)
                    text += ocr_text + "\n"
            return text if text.strip() else "No text found in PDF."
    except Exception as e:
        return f"PDF Extraction Error: {str(e)}"

def extract_text_from_docx(file_path):
    """Extract text from DOCX files."""
    try:
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text if text.strip() else "No text found in DOCX."
    except Exception as e:
        return f"DOCX Extraction Error: {str(e)}"

def extract_text_from_pptx(file_path):
    """Extract text from PPTX files."""
    try:
        prs = pptx.Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text if text.strip() else "No text found in PPTX."
    except Exception as e:
        return f"PPTX Extraction Error: {str(e)}"

def extract_data_from_csv(file_path):
    """Extract data from CSV files."""
    try:
        df = pd.read_csv(file_path)
        return df, df.to_string()
    except Exception as e:
        return None, f"CSV Extraction Error: {str(e)}"

def extract_data_from_xlsx(file_path):
    """Extract data from XLSX files."""
    try:
        df = pd.read_excel(file_path)
        return df, df.to_string()
    except Exception as e:
        return None, f"XLSX Extraction Error: {str(e)}"

def extract_data_from_json(file_path):
    """Extract data from JSON files."""
    try:
        with open(file_path, 'r') as f:
            data = json.load(f)
        return pd.json_normalize(data), json.dumps(data, indent=2)
    except Exception as e:
        return None, f"JSON Extraction Error: {str(e)}"

def crawl_url(url):
    """Crawl content from a URL."""
    try:
        response = requests.get(url, timeout=5)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, 'html.parser')
        text = ' '.join(p.get_text() for p in soup.find_all('p'))
        return text if text.strip() else "No content found on page."
    except Exception as e:
        return f"Web Crawling Error: {str(e)}"

def extract_hyperlinks(text):
    """Extract URLs from text."""
    url_pattern = re.compile(r'https?://[^\s<>"]+|www\.[^\s<>"]+')
    return url_pattern.findall(text)

def process_file(file_path):
    """Process a single file and extract content."""
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    df = None
    filename = os.path.basename(file_path)

    if ext == '.pdf':
        text = extract_text_from_pdf(file_path)
    elif ext == '.docx':
        text = extract_text_from_docx(file_path)
    elif ext == '.pptx':
        text = extract_text_from_pptx(file_path)
    elif ext in ['.png', '.jpg', '.jpeg']:
        text = extract_text_from_image(file_path)
    elif ext == '.csv':
        df, text = extract_data_from_csv(file_path)
    elif ext == '.xlsx':
        df, text = extract_data_from_xlsx(file_path)
    elif ext == '.json':
        df, text = extract_data_from_json(file_path)
    elif ext == '.txt':
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
    else:
        text = f"Unsupported file format: {ext}"

    if text and "Error" not in text:
        urls = extract_hyperlinks(text)
        for url in urls:
            crawled_text = crawl_url(url)
            text += f"\n\nContent from {url}:\n{crawled_text}"

    return text, df, filename